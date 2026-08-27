# -*- coding: utf-8 -*-
"""Genere docs/Recherche-existence-chemin.pptx

Resume illustre de la partie « recherche d'existence d'un chemin » de la
specification Word. Presentation par images (schemas dans docs/img/ppt/,
generes par docs/images_ppt.py) et par tableaux.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DOCS = Path(r"C:\Users\amami\GitHub\restitutiondonnees") / "docs"
IMG = DOCS / "img" / "ppt"
OUT = DOCS / "Recherche-existence-chemin.pptx"

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
INK = RGBColor(0x22, 0x2A, 0x35)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF1, 0xF3, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _box(slide, x, y, w, h):
    return slide.shapes.add_textbox(x, y, w, h)


def slide_header(slide, title, kicker=None):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.16))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tb = _box(slide, Inches(0.6), Inches(0.34), SW - Inches(1.2), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True
    if kicker:
        r = tf.paragraphs[0].add_run(); r.text = kicker.upper()
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = GREEN
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = NAVY


def new_slide(title=None, kicker=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        slide_header(s, title, kicker)
    return s


def bullets(slide, items, x=Inches(0.7), y=Inches(1.7), w=None, h=None, size=16, gap=7):
    w = w or (SW - Inches(1.4)); h = h or (SH - y - Inches(0.4))
    tf = _box(slide, x, y, w, h).text_frame; tf.word_wrap = True
    first = True
    for it in items:
        lvl, text = it if isinstance(it, tuple) else (0, it)
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.level = lvl; p.space_after = Pt(gap)
        r = p.add_run(); r.text = ("•  " if lvl == 0 else "–  ") + text
        r.font.size = Pt(size if lvl == 0 else size - 2)
        r.font.color.rgb = INK if lvl == 0 else GREY


def table(slide, rows, x=Inches(0.7), y=Inches(1.7), w=None, h=None,
          font=12, col_widths=None):
    w = w or (SW - Inches(1.4)); h = h or Inches(0.42 * len(rows))
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h).table
    if col_widths:
        tot = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(w * cw / tot))
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.margin_left = c.margin_right = Pt(7)
            c.margin_top = c.margin_bottom = Pt(3)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.text = str(val)
            for run in c.text_frame.paragraphs[0].runs:
                run.font.size = Pt(font)
                if ri == 0:
                    run.font.bold = True; run.font.color.rgb = WHITE
                else:
                    run.font.color.rgb = INK
            c.fill.solid()
            c.fill.fore_color.rgb = NAVY if ri == 0 else (WHITE if ri % 2 else LIGHT)
    return tbl


def code(slide, text, x=Inches(0.7), y=Inches(1.7), w=None, h=None, size=12):
    w = w or (SW - Inches(1.4)); h = h or (SH - y - Inches(0.4))
    box = slide.shapes.add_shape(1, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFA)
    box.line.color.rgb = RGBColor(0xD0, 0xD7, 0xDE); box.line.width = Pt(0.75)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(12); tf.margin_top = tf.margin_bottom = Pt(10)
    for i, line in enumerate(text.strip("\n").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line or " "
        r.font.name = "Consolas"; r.font.size = Pt(size); r.font.color.rgb = INK


def picture(slide, name, x, y, w=None, h=None):
    path = IMG / name
    if not path.exists():
        _box(slide, x, y, w or Inches(4), Inches(0.5)).text_frame.text = f"[{name} manquant]"
        return
    kw = {}
    if w: kw["width"] = w
    if h: kw["height"] = h
    slide.shapes.add_picture(str(path), x, y, **kw)


def caption(slide, text, x, y, w):
    tb = _box(slide, x, y, w, Inches(0.4))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = text
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r.font.size = Pt(10); r.font.italic = True; r.font.color.rgb = GREY


def footnote(slide, text):
    tb = _box(slide, Inches(0.7), SH - Inches(0.55), SW - Inches(1.4), Inches(0.45))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(10); r.font.italic = True; r.font.color.rgb = GREY


LX, LW = Inches(0.7), Inches(6.5)
RX, RW = Inches(7.5), Inches(5.2)
CY = Inches(1.8)

# ======================================================================
# 1 — Titre
# ======================================================================
s = new_slide()
band = s.shapes.add_shape(1, 0, Inches(2.1), SW, Inches(2.7))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
tb = _box(s, Inches(1.0), Inches(2.45), Inches(8.2), Inches(1.7))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run(); r.text = "Recherche d'existence d'un chemin"
r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = WHITE
r = tf.add_paragraph().add_run()
r.text = "Résumé illustré — spécification PathFinder (ASP.NET Core MVC)"
r.font.size = Pt(17); r.font.color.rgb = RGBColor(0xC7, 0xD2, 0xE0)
picture(s, "graphe-oriente.png", Inches(9.6), Inches(2.15), h=Inches(2.6))
tb = _box(s, Inches(1.0), Inches(5.2), Inches(11), Inches(0.6))
r = tb.text_frame.paragraphs[0].add_run()
r.text = "Graphe orienté non pondéré · BFS bidirectionnel · dbo.LINE_VIS_EDG · v3.0 — 2026-08-27"
r.font.size = Pt(12); r.font.color.rgb = GREY

# ======================================================================
# 2 — Le probleme
# ======================================================================
s = new_slide("Le problème à résoudre", "Objectif")
table(s, [
    ["Question", "Réponse attendue", "Parcours nécessaire"],
    ["Existence : un chemin orienté source → cible ?", "oui / non", "n'importe lequel (BFS ou DFS)"],
    ["Plus court chemin", "un chemin de longueur minimale (nb d'arêtes)", "par distance croissante → BFS"],
], y=Inches(1.9), font=13, col_widths=[4.5, 4, 4])
bullets(s, [
    "Résoudre le plus court chemin résout gratuitement l'existence : un seul BFS traite les deux.",
    "La volumétrie interdit de charger le graphe entier → parcours borné, requêtes ciblées.",
], y=Inches(4.1), size=15)

# ======================================================================
# 3 — Le graphe manipule
# ======================================================================
s = new_slide("Le graphe manipulé par l'application", "Contexte")
table(s, [
    ["Caractéristique", "RestitutionGraphe", "Conséquence"],
    ["Orienté", "oui (colonne Direction)", "le parcours respecte le sens"],
    ["Pondéré", "non", "BFS optimal ; Dijkstra inutile"],
    ["Acyclique (DAG)", "non garanti", "marquer les sommets visités"],
    ["Densité", "creux — 2 à 6 arêtes / nœud", "listes d'adjacence"],
    ["Taille", "~100 000 nœuds", "jamais tout le graphe"],
    ["Connexité", "non garantie", "beaucoup de couples sans chemin"],
], x=LX, y=CY, w=Inches(6.7), font=11.5, col_widths=[2.6, 3.2, 3.4])
picture(s, "graphe-oriente.png", Inches(8.0), Inches(2.4), h=Inches(3.6))
caption(s, "petit graphe orienté non pondéré", Inches(7.8), Inches(6.1), Inches(4.6))

# ======================================================================
# 4 — LINE_VIS_EDG
# ======================================================================
s = new_slide("La source de vérité : dbo.LINE_VIS_EDG", "Modèle de données")
table(s, [
    ["Colonne", "Type", "Rôle"],
    ["Nodes", "VARCHAR(8000)", "un nœud (extrémité 1)"],
    ["Direction", "chaîne", "'predecesseur' / 'successeur'"],
    ["NodesLie", "VARCHAR(8000)", "l'autre nœud (extrémité 2)"],
    ["Transformation", "nullable", "SELECT, JOIN… (affichée)"],
], x=LX, y=CY, w=Inches(6.6), font=12, col_widths=[3, 3, 4])
picture(s, "direction.png", Inches(7.6), Inches(1.9), w=Inches(5.4))
caption(s, "la colonne Direction encode le sens de l'arête", Inches(7.5), Inches(3.6), Inches(5.6))
table(s, [
    ["Pour un nœud X", "Lignes"],
    ["arêtes sortantes", "(Nodes=X, predecesseur)  ou  (NodesLie=X, successeur)"],
    ["arêtes entrantes", "(NodesLie=X, predecesseur)  ou  (Nodes=X, successeur)"],
], x=LX, y=Inches(4.5), w=Inches(11.9), font=12, col_widths=[3, 9])

# ======================================================================
# 5 — Pourquoi BFS
# ======================================================================
s = new_slide("Pourquoi BFS est le bon choix ici", "Choix de l'algorithme")
table(s, [
    ["Atout du contexte", "Ce que BFS en fait"],
    ["Graphe non pondéré", "rend le plus court chemin en O(n+m), sans file de priorité"],
    ["On veut juste l'existence", "s'arrête dès qu'il touche la cible"],
    ["« Aucun chemin » fréquents", "borné (maxDepth 20 ; 30 000 nœuds/sens) → réponse négative maîtrisée"],
    ["Données distantes (SQL)", "un palier = une requête par lot (pas une requête par arête)"],
    ["Symétrie source / cible", "se dédouble en BFS bidirectionnel → profondeur ÷ 2"],
], y=CY, font=12.5, col_widths=[3.5, 8.5])
footnote(s, "BFS est le seul algorithme qui donne le plus court chemin d'un graphe non pondéré "
            "au prix d'un simple parcours.")

# ======================================================================
# 6 — BFS bidirectionnel : principe
# ======================================================================
s = new_slide("BFS bidirectionnel — principe", "Algorithme")
table(s, [
    ["Élément", "Rôle"],
    ["front avant", "part de la source, suit les arêtes sortantes"],
    ["front arrière", "part de la cible, remonte les arêtes entrantes"],
    ["forwardPrev[X]", "prédécesseur de X sur le chemin depuis la source"],
    ["backwardNext[X]", "successeur de X sur le chemin vers la cible"],
    ["rencontre", "X connu des deux côtés → plus court chemin trouvé"],
], x=LX, y=CY, w=Inches(6.4), font=11.5, col_widths=[3, 5.5])
picture(s, "bfs-bidi.png", Inches(7.3), Inches(2.2), w=Inches(5.7))
caption(s, "les deux fronts avancent en alternance et se rejoignent au milieu",
        Inches(7.2), Inches(5.2), Inches(5.9))

# ======================================================================
# 7 — Deroule
# ======================================================================
s = new_slide("BFS bidirectionnel — déroulé", "Algorithme")
table(s, [
    ["#", "Étape"],
    ["1", "Vérifier l'existence de source et cible (sinon : non trouvé)."],
    ["2", "Si source = cible : renvoyer [source] (longueur 0)."],
    ["3", "Initialiser forwardPrev={source}, backwardNext={cible}, les deux frontières."],
    ["4", "Répéter ≤ maxDepth fois, en alternant : palier pair = front avant, impair = front arrière."],
    ["5", "Front avant : (u→v), u connu, v inconnu → forwardPrev[v]=u ; v ∈ backwardNext → rencontre."],
    ["6", "Front arrière : (u→v), v connu, u inconnu → backwardNext[u]=v ; u ∈ forwardPrev → rencontre."],
    ["7", "Frontières vides ou maxDepth atteint sans rencontre → non trouvé."],
], y=CY, font=12, col_widths=[0.7, 12])

# ======================================================================
# 8 — Pseudo-code
# ======================================================================
s = new_slide("Pseudo-code de référence", "Algorithme")
code(s, """
fonction PlusCourtChemin(source, cible, maxDepth <= 20):
    si non Existe(source) ou non Existe(cible): retourner NON_TROUVE
    si source == cible: retourner [source]

    forwardPrev  <- { source: NUL } ; forwardFrontier  <- [source]
    backwardNext <- { cible:  NUL } ; backwardFrontier <- [cible]

    pour step de 0 a maxDepth-1:
        si les deux frontieres sont vides: arreter
        si (step pair) et forwardFrontier et |forwardPrev| < 30000:
            pour chaque (u -> v) dans FetchEdgesFrom(forwardFrontier):
                si u inconnu (avant) ou v deja connu (avant): continuer
                forwardPrev[v] <- u
                si v dans backwardNext: retourner Recoller(v)
            forwardFrontier <- nouveaux v
        sinon si backwardFrontier et |backwardNext| < 30000:
            pour chaque (u -> v) dans FetchEdgesInto(backwardFrontier):
                si v inconnu (arriere) ou u deja connu (arriere): continuer
                backwardNext[u] <- v
                si u dans forwardPrev: retourner Recoller(u)
            backwardFrontier <- nouveaux u
    retourner NON_TROUVE
""", y=Inches(1.7), size=12)

# ======================================================================
# 9 — Exemple
# ======================================================================
s = new_slide("Exemple déroulé — recherche N1 → N4", "Algorithme")
picture(s, "exemple.png", Inches(7.7), Inches(1.9), w=Inches(5.3))
table(s, [
    ["Palier", "Front", "Frontière", "Rencontre ?"],
    ["0", "avant", "[N1] → N2, N9", "non"],
    ["1", "arrière", "[N4] → N3", "non"],
    ["2", "avant", "[N2, N9] → N3", "OUI (N3)"],
], x=LX, y=Inches(2.0), w=Inches(6.6), font=12, col_widths=[1.2, 1.4, 3, 1.8])
tb = _box(s, LX, Inches(4.3), Inches(6.6), Inches(0.8))
r = tb.text_frame.paragraphs[0].add_run()
r.text = "Chemin : N1 → N2 → N3 → N4   (longueur 3)"
r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = GREEN

# ======================================================================
# 10 — Complexite
# ======================================================================
s = new_slide("Complexité", "Algorithme")
picture(s, "complexite.png", Inches(7.4), Inches(1.9), w=Inches(5.6))
table(s, [
    ["Variante", "Nœuds visités"],
    ["BFS à sens unique", "≈ dᴿ"],
    ["BFS bidirectionnel", "≈ 2 · d^(R/2)"],
], x=LX, y=Inches(2.1), w=Inches(6.4), font=13, col_widths=[3.5, 3])
bullets(s, [
    "d = degré moyen, R = longueur du chemin.",
    "Coût réel = allers-retours SQL : ≤ 2·maxDepth lots, sous-requêtes de 1000 paramètres.",
    "Garde-fous : maxDepth ≤ 20 ; 30 000 nœuds / sens.",
], x=LX, y=Inches(3.7), w=Inches(6.4), size=13)

# ======================================================================
# 11 — BFS vs les autres
# ======================================================================
s = new_slide("BFS vs les autres algorithmes", "Comparaison")
table(s, [
    ["Algorithme", "Poids", "Complexité", "Plus court chemin ?", "Ici ?"],
    ["BFS / BFS bidirectionnel", "non", "O(n+m) ; ~2·d^(R/2)", "OUI (non pondéré)", "RETENU"],
    ["DFS", "n/a", "O(n+m)", "NON", "non"],
    ["Dijkstra", "≥ 0", "O((n+m) log n)", "OUI", "non (pas de poids)"],
    ["Bellman-Ford", "quelconques", "O(n·m)", "OUI", "non"],
    ["A*", "≥ 0 + heuristique", "≤ Dijkstra", "OUI", "non (pas d'heuristique)"],
    ["Floyd-Warshall", "quelconques", "O(n³)", "OUI (toutes paires)", "non (n trop grand)"],
    ["Union-Find", "n/a", "~O(m·α(n))", "NON", "en pré-calcul"],
    ["Tri topologique", "quelconques", "O(n+m)", "OUI si acyclique", "non (cycles)"],
], y=Inches(1.7), font=11.5, col_widths=[3.2, 2.2, 2.6, 2.6, 2.4])
footnote(s, "Les algos pondérés résolvent un problème plus dur à un coût supérieur ; "
            "DFS ne donne pas le plus court chemin ; le tri topologique exige l'acyclicité.")

# ======================================================================
# 12 — Regles de gestion
# ======================================================================
s = new_slide("Règles de gestion clés", "Spécification")
table(s, [
    ["Réf.", "Règle"],
    ["RG-01", "Un nœud existe s'il apparaît en colonne Nodes ou NodesLie."],
    ["RG-02", "source = cible (nœud existant) → chemin de longueur 0."],
    ["RG-03", "source ou cible inexistante → « aucun chemin », sans parcours."],
    ["RG-04", "Aucun chemin en ≤ maxDepth paliers → « aucun chemin »."],
    ["RG-07", "Plus court = nombre d'arêtes (les arêtes n'ont pas de poids)."],
    ["RG-08", "maxDepth plafonné à 20 (défaut 12)."],
    ["RG-09", "Résultat mis en cache 5 min ; cache borné à 10 000 entrées."],
    ["RG-10", "Chaque sens s'arrête au-delà de 30 000 nœuds visités."],
    ["RG-11", "Paramètres SQL typés VarChar(8000) → index préservés."],
    ["RG-12", "Requêtes sur liste de nœuds : lots de 1000 paramètres."],
], y=Inches(1.6), font=12, col_widths=[1.2, 11])

# ======================================================================
# 13 — Interface
# ======================================================================
s = new_slide("Interface — la route « / »", "Spécification")
table(s, [
    ["État de la page", "Affichage"],
    ["vierge", "le formulaire seul (arrivée sur /)"],
    ["erreur de saisie", "bandeau « Renseigne un nœud source ET un nœud cible »"],
    ["chemin trouvé", "bandeau vert + chaîne « N1 → N2 → … » + tableau (# / De / Vers / Transformation)"],
    ["chemin de longueur 0", "bandeau vert + note (source = cible)"],
    ["aucun chemin", "bandeau rouge"],
], y=CY, font=12, col_widths=[2.8, 9.2])
bullets(s, [
    "<form method=\"get\"> → la page se recharge sur /?source=…&target=… (URL partageable).",
    "Rendu 100 % côté serveur (Razor) — aucune API, aucun JavaScript. Résultat en cache 5 min.",
], y=Inches(5.1), size=13)

# ======================================================================
# 14 — Pistes d'amelioration
# ======================================================================
s = new_slide("Pistes d'amélioration de la base", "Évolutions")
table(s, [
    ["Piste", "Idée", "Effet sur l'existence", "État"],
    ["Index composites", "(Nodes,Direction) / (NodesLie,Direction) INCLUDE", "seek au lieu de scan", "recommandé"],
    ["Table EDGE normalisée", "pré-dériver Direction : EDGE(SourceId, TargetId)", "requêtes BFS triviales", "recommandé"],
    ["Pré-calcul des composantes", "1 scan → ComponentId par nœud (Union-Find)", "aucun chemin en O(1)", "IMPLÉMENTÉ (dotnet-new-scan/)"],
    ["Condensation SCC", "contracter les cycles → DAG + atteignabilité", "OUI/NON exact sans BFS", "IMPLÉMENTÉ (dotnet-new-scan/)"],
    ["Graphe en mémoire", "charger EDGE au démarrage (adjacence + inverse)", "latence SQL ÷ 10-100", "recommandé"],
], y=Inches(1.7), font=10.5, col_widths=[2.6, 4.6, 2.9, 2.4])
footnote(s, "Aucune ne change les règles de gestion. Accès données isolé dans "
            "LineVisEdgRepository → chaque piste est un changement localisé. "
            "Le pré-calcul des composantes est détaillé aux 2 slides suivantes.")

# ======================================================================
# 15 — Le scan : l'intuition des iles
# ======================================================================
s = new_slide("Le scan : l'intuition des îles", "Zoom (implémenté dans dotnet-new-scan/)")
picture(s, "composantes.png", Inches(7.2), Inches(3.0), w=Inches(5.7))
bullets(s, [
    "Le graphe = un archipel. Nœuds = maisons, arêtes = ponts.",
    "Un groupe de maisons reliées par des ponts = une île (composante connexe).",
    "Deux îles n'ont aucun pont entre elles.",
    "Donc : source et cible sur des îles différentes → AUCUN chemin, "
    "sans rien chercher.",
], x=LX, y=Inches(1.8), w=Inches(6.3), size=14)
table(s, [
    ["Île (ComponentId)", "Nœuds  (jeu de test)"],
    ["0", "N1 … N100000  (graphe d'origine)"],
    ["1", "X1, X2, X3, X4, X5"],
    ["2", "Y1, Y2, Y3"],
], x=LX, y=Inches(4.3), w=Inches(6.3), font=11)

# ======================================================================
# 16 — Le scan : Union-Find + utilisation
# ======================================================================
s = new_slide("Le scan : Union-Find, puis usage", "Zoom")
table(s, [
    ["Arête lue", "Îles après fusion"],
    ["(départ)", "{X1} {X2} {X3} {X4} {X5}"],
    ["X1—X2", "{X1,X2} {X3} {X4} {X5}"],
    ["X2—X3", "{X1,X2,X3} {X4} {X5}"],
    ["X3—X1", "{X1,X2,X3} …  (déjà ensemble)"],
    ["X3—X4", "{X1,X2,X3,X4} {X5}"],
    ["X4—X5", "{X1,X2,X3,X4,X5}  → île n° 1"],
], x=LX, y=Inches(1.8), w=Inches(6.6), font=11, col_widths=[2.2, 6])
bullets(s, [
    "1 seul balayage de LINE_VIS_EDG. Au départ chaque nœud est sa propre île ; "
    "chaque arête fusionne deux îles.",
    "Résultat écrit dans dbo.NODE_COMPONENT (NodeId → ComponentId).",
    "À chaque recherche :",
    (1, "île(source) ≠ île(cible) → « aucun chemin », O(1), SANS BFS ;"),
    (1, "île(source) = île(cible) → BFS comme d'habitude."),
    "Limite : prouve un NON, pas un OUI orienté (X5→X1 : même île, mais pas de "
    "chemin dans le sens → le BFS tranche).",
], x=Inches(7.6), y=Inches(1.8), w=Inches(5.3), size=12)
footnote(s, "Coût : O(n+m), une fois (~2,2 s pour 100 000 nœuds). Comparaison par "
            "recherche : O(1). À relancer via /Scan après modification de la table.")

# ======================================================================
# 17 — Le scan : le code
# ======================================================================
s = new_slide("Le scan : le code (dotnet-new-scan/)", "Zoom")
code(s, """
// Union-Find  —  GraphScanService.cs
public string Find(string x) {
    Ensure(x);                    // 1er contact : x racine de lui-meme
    while (_parent[x] != x) {
        _parent[x] = _parent[_parent[x]];   // compression de chemin
        x = _parent[x];
    }
    return x;
}
public void Union(string a, string b) {
    var ra = Find(a); var rb = Find(b);
    if (ra == rb) return;
    if (_rank[ra] < _rank[rb]) (ra, rb) = (rb, ra);  // union par rang
    _parent[rb] = ra;
    if (_rank[ra] == _rank[rb]) _rank[ra]++;
}
""", x=LX, y=Inches(1.7), w=Inches(6.5), h=Inches(4.3), size=11)
code(s, """
Service : PAS de SQL. Repositories : tout le SQL.

GraphScanService.Run() :
  1. foreach (u,v) in _edges.StreamAllEdges()
         uf.Union(u, v)
  2. foreach node : root = uf.Find(node)
         root nouveau -> id 0,1,2,...
         rows.Add((node, id))
  3. _components.ReplaceAll(rows)      // repo

GraphScanService.Compare(src, dst) :
  (cs, ct) = _components.GetComponentIds(src, dst)
  cs/ct null  -> ScanUnavailable  (BFS)
  cs != ct    -> DifferentComponents ("aucun chemin", O(1))
  cs == ct    -> SameComponent (BFS)

  LineVisEdgRepository.StreamAllEdges()  -> SELECT Nodes,NodesLie
  NodeComponentRepository.ReplaceAll()   -> DROP/CREATE/BulkCopy
  NodeComponentRepository.GetComponentIds() -> SELECT ... IN (@a,@b)
""", x=Inches(7.5), y=Inches(1.7), w=Inches(5.3), h=Inches(4.5), size=9.5)
footnote(s, "Fichiers : Services/GraphScanService.cs (algorithme) · "
            "Models/LineVisEdgRepository.cs + Models/NodeComponentRepository.cs (SQL).")

# ======================================================================
# 18 — Condensation SCC (§ 11.5)
# ======================================================================
s = new_slide("Aller plus loin : la condensation SCC", "Zoom (§ 11.5, dotnet-new-scan/)")
bullets(s, [
    "Composante FORTEMENT connexe (SCC) : nœuds tous reliés par un chemin "
    "orienté dans les deux sens.",
    "On contracte chaque SCC en un super-nœud → le « graphe condensé », "
    "toujours un DAG, beaucoup plus petit.",
    "Chemin orienté u → v existe  ⟺  SccId(v) atteignable depuis SccId(u) "
    "dans le DAG condensé (BFS sur ~2000 super-nœuds, pas 100 000).",
    "Verdict EXACT : oui comme non, sans BFS sur le graphe d'origine.",
    "Algorithme : Kosaraju (2 DFS itératifs), O(n+m).",
], x=LX, y=Inches(1.8), w=Inches(6.4), size=13)
table(s, [
    ["", "Graphe d'origine", "Graphe condensé"],
    ["Nœuds", "100 000", "1 979 SCC"],
    ["Arêtes", "~400 000", "2 144"],
    ["Structure", "cyclique", "DAG"],
    ["+ grande SCC", "—", "98 028 nœuds"],
], x=Inches(7.4), y=Inches(2.0), w=Inches(5.4), font=11, col_widths=[2, 3, 3])
footnote(s, "Gain sur le § 11.4 : X5 → X1 (même île, mais pas de chemin dans le "
            "sens) → « aucun chemin » exact et immédiat. Le § 11.4 renvoyait au BFS.")

# ======================================================================
# 19 — Synthese
# ======================================================================
s = new_slide("Synthèse", None)
table(s, [
    ["Aspect", "En bref"],
    ["Le graphe", "orienté, non pondéré, creux, cyclique possible, non connexe, ~100 000 nœuds"],
    ["Le besoin", "existence + plus court chemin — résolus ensemble par un seul BFS"],
    ["L'algorithme", "BFS bidirectionnel par paliers, exécuté en SQL, borné, résultat en cache 5 min"],
    ["Pas Dijkstra / A* / Floyd-Warshall", "problème plus dur (poids, toutes paires), coût supérieur, sans bénéfice ici"],
    ["Pré-calculs (dotnet-new-scan/)", "§ 11.4 composantes faibles (NON certain) · § 11.5 condensation SCC (OUI/NON exact sans BFS)"],
], y=Inches(1.9), font=12, col_widths=[3.4, 8.6])

prs.save(OUT)
print(f"OK -> {OUT}  ({len(prs.slides)} diapositives)")
